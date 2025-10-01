import os
from typing import Tuple

def sniff_mime(filename: str) -> str:
    f = filename.lower()
    if f.endswith(".docx"):
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if f.endswith(".xlsx"):
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if f.endswith(".pptx"):
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if f.endswith(".txt"):
        return "text/plain"
    return "application/octet-stream"

def extract_excerpt(filepath: str, mime: str, limit: int = 4000) -> str:
    try:
        if mime.endswith("wordprocessingml.document"):
            from docx import Document
            doc = Document(filepath)
            text = "\n".join(p.text for p in doc.paragraphs)
            return text[:limit]

        if mime.endswith("spreadsheetml.sheet"):
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets[:3]:
                for row in ws.iter_rows(min_row=1, max_row=50, values_only=True):
                    line = " ".join("" if c is None else str(c) for c in row)
                    texts.append(line)
                    if sum(len(t) for t in texts) > limit:
                        break
                if sum(len(t) for t in texts) > limit:
                    break
            return "\n".join(texts)[:limit]

        if mime.endswith("presentationml.presentation"):
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = []
            for slide in prs.slides[:10]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
                if sum(len(t) for t in texts) > limit:
                    break
            return "\n".join(texts)[:limit]

        if mime == "text/plain":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(limit)

        with open(filepath, "rb") as f:
            data = f.read(limit)
        return data.decode("utf-8", errors="ignore")
    except Exception as e:
        return f"[extract_error] {type(e).__name__}: {e}"
